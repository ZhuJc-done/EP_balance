"""Generate the two-slide editable architecture deck requested for Scale-EPLB."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "arch.pptx"

WHITE = "FFFFFF"
BG = "F7F9FC"
INK = "172033"
MUTED = "5B657A"
BORDER = "334155"
GRID = "CBD5E1"
LIGHT_GRAY = "EEF2F6"
ORANGE = "E56A2E"
ORANGE_LIGHT = "FFF1E8"
GREEN = "2E7D32"
GREEN_LIGHT = "EAF5E8"
BLUE = "2F67B1"
BLUE_LIGHT = "EAF2FC"
PURPLE = "7257A8"
PURPLE_LIGHT = "F1ECF8"
RED = "B03030"
RED_LIGHT = "F7DEDE"
L2_BLUE = "2F6FDB"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_shape_text(
    shape,
    value: str,
    *,
    size: float = 10,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    font: str = "Aptos",
    italic: bool = False,
    margin: float = 0.04,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for index, line_value in enumerate(value.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = line_value
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(color)


def box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = BORDER,
    width: float = 1,
    rounded: bool = True,
    text_value: str = "",
    size: float = 10,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    font: str = "Aptos",
    italic: bool = False,
    dash: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    if dash:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if text_value:
        set_shape_text(
            shape,
            text_value,
            size=size,
            color=color,
            bold=bold,
            align=align,
            font=font,
            italic=italic,
        )
    return shape


def textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    value: str,
    *,
    size: float = 10,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Aptos",
    italic: bool = False,
):
    shape = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_shape_text(
        shape,
        value,
        size=size,
        color=color,
        bold=bold,
        align=align,
        font=font,
        italic=italic,
        margin=0,
    )
    return shape


def add_arrowhead(connector, color: str) -> None:
    line = connector._element.spPr.ln
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    line.append(tail)
    connector.line.color.rgb = rgb(color)


def connector(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = BORDER,
    width: float = 1.4,
    arrow: bool = True,
    dash: bool = False,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        add_arrowhead(line, color)
    return line


def title(slide, main: str, subtitle: str) -> None:
    textbox(
        slide,
        0.25,
        0.08,
        12.83,
        0.38,
        main,
        size=22,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        0.25,
        0.46,
        12.83,
        0.22,
        subtitle,
        size=9.5,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_pipeline_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)
    title(
        slide,
        "Scale-EPLB GPU Parallel Solver Pipeline",
        "Two logical algorithm stages · Three ordered CUDA kernels",
    )

    # Input panel
    box(slide, 0.18, 0.78, 1.43, 2.68, fill=WHITE, line=BORDER, rounded=True)
    textbox(slide, 0.28, 0.84, 1.23, 0.24, "INPUT", size=11, bold=True, align=PP_ALIGN.CENTER)
    for i, label in enumerate(
        ("Load Ω [R,E]", "Topology / Domains", "Main + Slot Budget", "Integer Config")
    ):
        box(
            slide,
            0.29,
            1.18 + i * 0.36,
            1.21,
            0.27,
            fill=WHITE,
            line=MUTED,
            rounded=True,
            text_value=label,
            size=7.5,
            bold=i == 0,
        )
    connector(slide, 0.90, 2.52, 0.90, 2.72, color=PURPLE)
    box(
        slide,
        0.36,
        2.73,
        1.08,
        0.53,
        fill=PURPLE,
        line="4B3775",
        text_value="GPU\nGlobal Memory",
        size=9,
        color=WHITE,
        bold=True,
    )
    textbox(
        slide,
        0.27,
        3.25,
        1.25,
        0.16,
        "all tensors device-resident",
        size=6.6,
        color=MUTED,
        italic=True,
        align=PP_ALIGN.CENTER,
    )

    # Stage 1
    box(slide, 1.82, 0.78, 2.55, 2.68, fill=ORANGE_LIGHT, line=ORANGE)
    textbox(
        slide,
        1.94,
        0.84,
        2.31,
        0.22,
        "STAGE 1 · Cross-Domain Replica Admission",
        size=9.3,
        color=ORANGE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        2.08,
        1.07,
        2.03,
        0.18,
        "benefit-sorted (expert, domain) candidates",
        size=7,
        color=MUTED,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    box(
        slide,
        2.02,
        1.33,
        2.15,
        1.18,
        fill=WHITE,
        line=ORANGE,
        dash=True,
        text_value="",
    )
    textbox(
        slide,
        2.16,
        1.43,
        1.87,
        0.20,
        "Kernel 1 · Control (1 CUDA Block)",
        size=8.3,
        color=ORANGE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        2.15,
        1.71,
        1.30,
        0.48,
        "Process candidates\ndeterministically\n(fixed tie-breaks)",
        size=7.2,
        align=PP_ALIGN.CENTER,
    )
    for row in range(3):
        for col in range(2):
            box(
                slide,
                3.55 + col * 0.19,
                1.72 + row * 0.19,
                0.14,
                0.14,
                fill=ORANGE_LIGHT,
                line=ORANGE,
                rounded=False,
            )
    connector(slide, 2.62, 2.51, 2.62, 2.72, color=INK)
    connector(slide, 3.58, 2.51, 3.58, 2.72, color=INK)
    box(
        slide,
        2.11,
        2.74,
        1.02,
        0.42,
        fill=WHITE,
        line=ORANGE,
        text_value="x [E,R]\nint8",
        size=8,
        font="Aptos Mono",
    )
    box(
        slide,
        3.17,
        2.74,
        1.02,
        0.42,
        fill=WHITE,
        line=ORANGE,
        text_value="slot_used [R]\nint64",
        size=7.5,
        font="Aptos Mono",
    )
    textbox(
        slide,
        2.20,
        3.20,
        1.78,
        0.17,
        "actual domains M, not padded R",
        size=7,
        color=ORANGE,
        bold=True,
        italic=True,
        align=PP_ALIGN.CENTER,
    )

    # Route stage (dominant)
    box(slide, 4.58, 0.78, 4.49, 2.68, fill=GREEN_LIGHT, line=GREEN)
    textbox(
        slide,
        4.73,
        0.84,
        4.19,
        0.23,
        "ROUTE · Parallel Initial Routing / 并行初始路由",
        size=10.3,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        4.95,
        1.08,
        3.76,
        0.18,
        "Kernel 2 · R × E CUDA blocks distributed across SMs",
        size=7.5,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    textbox(slide, 4.76, 1.33, 1.30, 0.16, "R × E block grid", size=7.2, bold=True, align=PP_ALIGN.CENTER)
    for row in range(3):
        for col in range(4):
            box(
                slide,
                4.78 + col * 0.31,
                1.53 + row * 0.29,
                0.25,
                0.22,
                fill=WHITE,
                line=GREEN,
                rounded=False,
                text_value="src,e",
                size=5.2,
                font="Aptos Mono",
            )
    textbox(slide, 5.97, 1.78, 0.22, 0.22, "…", size=14, bold=True, align=PP_ALIGN.CENTER)
    for col, label in enumerate(("SM 0", "SM 1", "SM 2", "SM N−1")):
        x = 4.76 + col * 0.36
        box(slide, x, 2.52, 0.31, 0.52, fill=WHITE, line=GREEN, rounded=False)
        textbox(slide, x, 2.54, 0.31, 0.11, label, size=4.7, bold=True, align=PP_ALIGN.CENTER)
        for r in range(2):
            for c in range(2):
                box(
                    slide,
                    x + 0.04 + c * 0.10,
                    2.70 + r * 0.10,
                    0.07,
                    0.07,
                    fill="75C043",
                    line=GREEN,
                    rounded=False,
                )
    connector(slide, 5.30, 2.42, 5.30, 2.51, color=GREEN, dash=True)
    connector(slide, 5.84, 2.42, 5.84, 2.51, color=GREEN, dash=True)

    box(slide, 6.36, 1.34, 1.62, 1.73, fill=WHITE, line=GREEN)
    textbox(slide, 6.48, 1.42, 1.38, 0.20, "One (src, expert) Block", size=7.4, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 6.52, 1.66, 1.30, 0.15, "thread = destination rank", size=6.2, align=PP_ALIGN.CENTER)
    for i in range(8):
        box(
            slide,
            6.48 + i * 0.16,
            1.86,
            0.13,
            0.16,
            fill=GREEN_LIGHT if i < 7 else LIGHT_GRAY,
            line=GREEN,
            rounded=False,
            text_value=str(i) if i < 4 else ("R−1" if i == 7 else ""),
            size=5.2,
        )
    box(slide, 6.50, 2.10, 1.34, 0.22, fill=LIGHT_GRAY, line=GREEN, rounded=False, text_value="Warp 0   (32 threads)", size=6.2)
    box(slide, 6.50, 2.34, 1.34, 0.22, fill=LIGHT_GRAY, line=GREEN, rounded=False, text_value="Warp 1 … Warp W−1", size=6.0)
    box(
        slide,
        6.50,
        2.65,
        1.34,
        0.30,
        fill=GREEN_LIGHT,
        line=GREEN,
        text_value="Shared-memory prefix scan\nO(log R)",
        size=6.6,
        bold=True,
    )
    textbox(
        slide,
        8.05,
        1.53,
        0.84,
        1.38,
        "coalesced\nq writes\n\nunique q writer\n\nint64 atomicAdd\n→ U / rank_load",
        size=6.4,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.LEFT,
    )

    # Stage 2
    box(slide, 9.28, 0.78, 3.87, 2.68, fill=BLUE_LIGHT, line=BLUE)
    textbox(
        slide,
        9.42,
        0.84,
        3.59,
        0.22,
        "STAGE 2 · Intra-Domain Greedy Repair",
        size=9.5,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        9.76,
        1.07,
        2.92,
        0.17,
        "域内贪心修正",
        size=7.1,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    box(slide, 9.51, 1.31, 3.42, 1.86, fill=WHITE, line=BLUE, dash=True)
    textbox(slide, 9.68, 1.39, 3.08, 0.19, "Kernel 3 · Repair (1 CUDA Block)", size=8.2, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 10.21, 1.62, 2.02, 0.16, "up to 1024 threads", size=6.6, align=PP_ALIGN.CENTER)
    for i in range(14):
        box(
            slide,
            9.72 + i * 0.20,
            1.83,
            0.12,
            0.12,
            fill=BLUE_LIGHT,
            line=BLUE,
            rounded=False,
        )
    textbox(slide, 10.03, 2.00, 2.38, 0.15, "parallel expert / target scan", size=6.3, align=PP_ALIGN.CENTER)
    connector(slide, 11.22, 2.16, 11.22, 2.27, color=BLUE)
    box(
        slide,
        9.84,
        2.29,
        2.76,
        0.39,
        fill=BLUE_LIGHT,
        line=BLUE,
        text_value="Shared-memory deterministic reduction\nΔ↓ · load↑ · cost↑ · expert↑ · target↑",
        size=6.4,
        bold=True,
    )
    box(
        slide,
        9.84,
        2.74,
        2.76,
        0.22,
        fill=WHITE,
        line=BLUE,
        text_value="thread 0: floor-safe quota commit",
        size=6.5,
        italic=True,
    )
    textbox(slide, 9.84, 2.98, 1.36, 0.13, "≤ I_max repair rounds", size=6.3, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    box(slide, 11.32, 2.93, 1.29, 0.29, fill=BLUE_LIGHT, line=BLUE, text_value="Plan {x, q, θ}", size=7.5, bold=True, font="Aptos Mono")

    # Inter-stage arrows and synchronization line
    connector(slide, 1.61, 2.05, 1.81, 2.05, color=PURPLE, width=2)
    connector(slide, 4.37, 2.05, 4.57, 2.05, color=PURPLE, width=2)
    connector(slide, 9.07, 2.05, 9.27, 2.05, color=PURPLE, width=2)
    connector(slide, 0.18, 3.66, 13.15, 3.66, color=PURPLE, width=0.8, arrow=False, dash=True)
    textbox(
        slide,
        4.46,
        3.54,
        4.42,
        0.18,
        "same CUDA stream · kernel boundary = global synchronization",
        size=6.2,
        color=PURPLE,
        italic=True,
        align=PP_ALIGN.CENTER,
    )

    # Hardware execution mapping
    textbox(
        slide,
        0.18,
        3.80,
        12.97,
        0.22,
        "CUDA Execution Mapping Across Many SMs",
        size=10.2,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    box(slide, 0.18, 4.10, 1.08, 1.45, fill=WHITE, line=GRID, rounded=False)
    legend = (("75C043", "Route blocks"), (ORANGE, "Stage 1"), (BLUE, "Stage 2"))
    for i, (color, label) in enumerate(legend):
        box(slide, 0.29, 4.29 + i * 0.32, 0.10, 0.10, fill=color, line=color, rounded=False)
        textbox(slide, 0.44, 4.22 + i * 0.32, 0.70, 0.24, label, size=6.0, align=PP_ALIGN.LEFT)

    sm_x = [1.38, 3.30, 5.22, 7.14, 9.06, 11.00]
    sm_names = ["SM 0", "SM 1", "SM 2", "SM 3", "SM N−2", "SM N−1"]
    for x, name in zip(sm_x, sm_names):
        box(slide, x, 4.10, 1.73, 1.45, fill=WHITE, line=BORDER, rounded=False)
        textbox(slide, x, 4.14, 1.73, 0.17, name, size=7.2, bold=True, align=PP_ALIGN.CENTER)
        box(slide, x + 0.10, 4.37, 1.53, 0.18, fill=LIGHT_GRAY, line=GRID, rounded=False, text_value="Warp Schedulers", size=5.7)
        textbox(slide, x + 0.10, 4.60, 0.34, 0.14, "Warps", size=5.5, bold=True)
        for i in range(8):
            box(slide, x + 0.47 + i * 0.13, 4.61, 0.09, 0.09, fill="75C043", line=GREEN, rounded=False)
        textbox(slide, x + 0.10, 4.80, 0.34, 0.14, "Threads", size=5.5, bold=True)
        for i in range(8):
            box(slide, x + 0.47 + i * 0.13, 4.81, 0.09, 0.09, fill=WHITE, line=MUTED, rounded=False)
        box(slide, x + 0.10, 5.01, 1.53, 0.18, fill=ORANGE_LIGHT, line=ORANGE, rounded=False, text_value="Registers", size=5.8, bold=True)
        box(slide, x + 0.10, 5.24, 1.53, 0.20, fill=BLUE_LIGHT, line=BLUE, rounded=False, text_value="L1 / Shared Memory", size=5.7, bold=True)

    box(slide, 0.18, 5.69, 12.97, 0.27, fill=L2_BLUE, line="1E4FA0", rounded=False, text_value="L2 Cache · unified across all SMs", size=8.5, color=WHITE, bold=True)
    box(slide, 0.18, 6.06, 12.97, 0.31, fill=PURPLE, line="4B3775", rounded=False, text_value="GPU Global Memory", size=9.2, color=WHITE, bold=True)

    buffer_specs = (
        (1.37, 2.38, "q  int64 [R,E,R]\ndispatch plan · largest buffer"),
        (3.88, 1.67, "x  int8 [E,R]\nreplica decisions"),
        (5.68, 1.67, "U  int64 [E,R]\nper-expert usage"),
        (7.48, 1.76, "rank_load  int64 [R]\nper-rank load"),
    )
    textbox(slide, 0.18, 6.57, 1.03, 0.44, "Global buffers\n(native tensors)", size=6.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    for x, w, label in buffer_specs:
        box(slide, x, 6.52, w, 0.52, fill=WHITE, line=PURPLE, text_value=label, size=6.8, font="Aptos Mono")
    box(slide, 9.40, 6.51, 1.83, 0.53, fill="FFF9E8", line="A66A00", text_value="R=32, E=640\n92.64 ms → 1.17 ms\n~79× kernel speedup", size=6.8, bold=True)
    box(slide, 11.38, 6.51, 1.77, 0.53, fill=GREEN_LIGHT, line=GREEN, text_value="Integer decisions\nfixed tie-breaks\nno CPU sync", size=6.7, bold=True)
    textbox(slide, 0.18, 7.20, 12.97, 0.15, "R = ranks · E = experts · M = actual domains · all decision tensors use integer arithmetic", size=5.8, color=MUTED, italic=True)


def add_gin_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)
    title(
        slide,
        "GIN Replica Weight & Gradient Transport",
        "Device-initiated symmetric-window data path · no CPU proxy on the critical path",
    )

    # GPU A
    box(slide, 0.42, 0.84, 4.34, 3.30, fill=WHITE, line=BORDER)
    textbox(slide, 0.62, 0.92, 3.94, 0.25, "GPU A  ·  owner of main(e)", size=12, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 0.70, 1.19, 1.70, 0.18, "symmetric windows", size=7.5, color=MUTED, italic=True)
    box(slide, 0.77, 1.43, 3.64, 0.55, fill="CFE0F5", line="5A7FB0", text_value="home[j]\nWₑ  (main weights)", size=8.5, font="Aptos Mono")
    box(slide, 0.77, 2.19, 3.64, 0.55, fill="CDECCD", line="5AA05A", text_value="scratch[j]\ngrad columns  (local_slot × src)", size=8.2, font="Aptos Mono")
    box(slide, 0.77, 2.95, 3.64, 0.55, fill="F6DCB0", line="C79437", text_value="slot[j]\nlocal instances", size=8.5, font="Aptos Mono")
    textbox(slide, 0.75, 3.69, 3.68, 0.24, "(3) main reduces scratch columns → ∇Wₑ", size=8.5, bold=True, align=PP_ALIGN.CENTER)
    connector(slide, 0.72, 2.47, 0.72, 1.71, color=BORDER, width=1.5)

    # GPU B
    box(slide, 8.57, 0.84, 4.34, 3.30, fill=WHITE, line=BORDER)
    textbox(slide, 8.77, 0.92, 3.94, 0.25, "GPU B  ·  hosts replica of e", size=12, bold=True, align=PP_ALIGN.CENTER)
    box(slide, 8.92, 1.43, 3.64, 0.55, fill="E6E6E6", line="888888", text_value="device descriptor\npeers / off / nbytes  (device-resident)", size=8.2)
    box(slide, 8.92, 2.19, 3.64, 0.55, fill="F6DCB0", line="C79437", text_value="slot[j]\nreplica of e @ slot s", size=8.5, font="Aptos Mono")
    box(slide, 8.92, 2.95, 3.64, 0.55, fill="CDECCD", line="5AA05A", text_value="scratch[j]\nremote gradient staging", size=8.5, font="Aptos Mono")
    connector(slide, 10.74, 1.99, 10.74, 2.18, color=MUTED, width=1.0, dash=True)
    textbox(slide, 10.87, 1.98, 1.52, 0.17, "descriptor drives kernel", size=6.2, color=MUTED)

    # FWD/BWD paths
    connector(slide, 4.41, 1.70, 8.92, 2.42, color=BORDER, width=2.0)
    box(slide, 5.45, 1.42, 2.40, 0.33, fill=WHITE, line=WHITE, rounded=False, text_value="(1) FWD: GIN get", size=10, bold=True)
    textbox(slide, 5.87, 1.75, 1.56, 0.18, "GPU B pulls Wₑ", size=7.2, color=MUTED, align=PP_ALIGN.CENTER)
    connector(slide, 8.92, 2.68, 4.41, 2.45, color=BORDER, width=2.0)
    box(slide, 5.27, 2.64, 2.79, 0.33, fill=WHITE, line=WHITE, rounded=False, text_value="(2) BWD: GIN put ∇", size=10, bold=True)

    # Transport band
    connector(slide, 2.58, 4.14, 2.58, 4.37, color=MUTED, width=1.0)
    connector(slide, 10.73, 4.14, 10.73, 4.37, color=MUTED, width=1.0)
    box(
        slide,
        1.32,
        4.39,
        10.69,
        0.52,
        fill="D8EFCF",
        line="3C7D3C",
        text_value="GIN primitive  ⇒  RDMA / NVLink\nDevice-initiated GDAKI / IBGDA · no CPU on the critical path",
        size=9,
        color=GREEN,
        bold=True,
    )

    # Call stack
    textbox(
        slide,
        0.43,
        5.08,
        7.25,
        0.22,
        "Kernel-initiated RDMA path  ·  zoom-in of transport",
        size=9,
        bold=True,
    )
    labels = (
        "GPU kernel\nncclGin::get / put",
        "post WQE (QP)\nring doorbell",
        "NIC",
        "GPUDirect\nRDMA",
        "remote\nsymmetric window",
        "CQE / flush\n+ signal fence",
    )
    call_x = [0.42, 2.57, 4.72, 6.87, 9.02, 11.17]
    for index, (x, label) in enumerate(zip(call_x, labels)):
        box(slide, x, 5.39, 1.72, 0.65, fill=WHITE, line=MUTED, text_value=label, size=7.5, font="Aptos Mono" if index == 0 else "Aptos")
        if index < len(labels) - 1:
            connector(slide, x + 1.72, 5.72, call_x[index + 1], 5.72, color=MUTED, width=1.2)

    # Crossed-out CPU path
    box(slide, 4.73, 6.42, 1.72, 0.56, fill=RED_LIGHT, line=RED, text_value="CPU proxy", size=9, color=RED, bold=True)
    connector(slide, 4.73, 6.42, 6.45, 6.98, color=RED, width=1.4, arrow=False)
    connector(slide, 4.73, 6.98, 6.45, 6.42, color=RED, width=1.4, arrow=False)
    textbox(slide, 6.62, 6.57, 1.20, 0.22, "bypassed", size=8, color=RED, bold=True)
    connector(slide, 3.43, 6.04, 5.14, 6.41, color=RED, width=1.1, dash=True)
    textbox(
        slide,
        8.18,
        6.35,
        4.73,
        0.64,
        "Forward: replica directly pulls main weights\nBackward: replica directly pushes gradients\nOwner GPU reduces scratch columns into ∇Wₑ",
        size=7.6,
        color=MUTED,
        align=PP_ALIGN.LEFT,
    )
    textbox(slide, 0.42, 7.22, 12.49, 0.14, "All boxes, labels, arrows, memory windows, and call-stack elements are editable PowerPoint shapes.", size=6.2, color=MUTED, italic=True)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    add_pipeline_slide(prs)
    add_gin_slide(prs)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build())
